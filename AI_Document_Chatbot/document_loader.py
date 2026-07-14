"""
document_loader.py

Loads PDF and PowerPoint documents and extracts text.
"""

import os
import fitz
from pptx import Presentation


class DocumentLoader:

    def __init__(self, filepath):

        self.filepath = filepath

        self.filename = os.path.basename(filepath)

        self.extension = os.path.splitext(filepath)[1].lower()

    # --------------------------------------------------------

    def load_document(self):

        if self.extension == ".pdf":

            return self.read_pdf()

        elif self.extension in [".ppt", ".pptx"]:

            return self.read_ppt()

        else:

            raise ValueError(
                f"Unsupported file type: {self.extension}"
            )

    # --------------------------------------------------------

    @staticmethod
    def clean_text(text):

        if not text:

            return ""

        text = text.replace("\n", " ")

        text = " ".join(text.split())

        return text

    # --------------------------------------------------------

    def read_pdf(self):

        pages = []

        try:

            with fitz.open(self.filepath) as document:

                for page_number, page in enumerate(document, start=1):

                    text = page.get_text("text")

                    text = self.clean_text(text)

                    if text:

                        pages.append({

                            "page": page_number,

                            "text": text,

                            "source": self.filename

                        })

            return pages

        except Exception as e:

            raise Exception(
                f"Error reading PDF '{self.filename}': {e}"
            )

    # --------------------------------------------------------

    def read_ppt(self):

        slides = []

        try:

            presentation = Presentation(self.filepath)

            for slide_number, slide in enumerate(
                presentation.slides,
                start=1
            ):

                slide_text = []

                for shape in slide.shapes:

                    if hasattr(shape, "text"):

                        text = self.clean_text(shape.text)

                        if text:

                            slide_text.append(text)

                combined_text = " ".join(slide_text)

                if combined_text:

                    slides.append({

                        "slide": slide_number,

                        "text": combined_text,

                        "source": self.filename

                    })

            return slides

        except Exception as e:

            raise Exception(
                f"Error reading PowerPoint '{self.filename}': {e}"
            )