import os
from PyPDF2 import PdfReader
from pptx import Presentation


# -----------------------------------
# Read PDF File
# -----------------------------------

def read_pdf(file_path):

    text = ""

    try:

        reader = PdfReader(file_path)

        for page in reader.pages:

            page_text = page.extract_text()

            if page_text:
                text += page_text + "\n"

    except Exception as e:

        print("PDF Error:", e)

    return text


# -----------------------------------
# Read PowerPoint File
# -----------------------------------

def read_ppt(file_path):

    text = ""

    try:

        presentation = Presentation(file_path)

        for slide in presentation.slides:

            for shape in slide.shapes:

                if hasattr(shape, "text"):

                    text += shape.text + "\n"

    except Exception as e:

        print("PPT Error:", e)

    return text


# -----------------------------------
# Detect File Type
# -----------------------------------

def extract_text(file_path):

    extension = os.path.splitext(file_path)[1].lower()

    if extension == ".pdf":

        return read_pdf(file_path)

    elif extension in [".ppt", ".pptx"]:

        return read_ppt(file_path)

    else:

        return ""